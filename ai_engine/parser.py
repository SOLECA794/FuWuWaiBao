import fitz  # PyMuPDF
import os
import json
import time
import re
from typing import Any
from pathlib import Path
from pptx import Presentation

try:
    from .schema import build_document_schema
except ImportError:
    from schema import build_document_schema

class DocumentParser:
    """AI工程师：文档解析模块，负责将PDF/PPT转化为带页码的结构化文本"""
    
    def __init__(self, file_path):
        self.file_path = file_path
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件未找到: {file_path}")

    def parse_pdf_pages(self) -> list[dict[str, Any]]:
        """
        解析PDF文件，返回页级结构（兼容旧逻辑）
        返回格式: [{"page": 1, "content": "..."}, ...]
        """
        doc = fitz.open(self.file_path)
        parsed_pages: list[dict[str, Any]] = []

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = self._extract_pdf_text(page)

            if self._is_meaningful_text(text):
                cleaned = self._clean_text(text)
                parsed_pages.append(
                    {
                        "page": page_num + 1,
                        "content": cleaned,
                        "content_length": len(cleaned),
                    }
                )

        doc.close()
        return parsed_pages

    def parse_pdf(self) -> dict[str, Any]:
        """
        解析PDF并返回统一协议的结构化JSON数据。
        """
        started_at = time.time()
        with fitz.open(self.file_path) as doc:
            total_pages = len(doc)

        parsed_pages = self.parse_pdf_pages()
        return build_document_schema(
            file_path=self.file_path,
            document_type="pdf",
            total_pages=total_pages,
            parsed_pages=parsed_pages,
            started_at=started_at,
        )

    def parse_pptx_pages(self) -> list[dict[str, Any]]:
        """
        解析PPTX文件，返回页级结构（每一页对应一张幻灯片）。
        返回格式: [{"page": 1, "content": "..."}, ...]
        """
        presentation = Presentation(self.file_path)
        parsed_pages: list[dict[str, Any]] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)

            notes_frame = getattr(getattr(slide, "notes_slide", None), "notes_text_frame", None)
            if notes_frame and getattr(notes_frame, "text", ""):
                texts.append(notes_frame.text)

            merged_text = "\n".join(texts).strip()
            if self._is_meaningful_text(merged_text):
                cleaned = self._clean_text(merged_text)
                parsed_pages.append(
                    {
                        "page": slide_index,
                        "content": cleaned,
                        "content_length": len(cleaned),
                    }
                )

        return parsed_pages

    def parse_pptx(self) -> dict[str, Any]:
        """
        解析PPTX并返回统一协议的结构化JSON数据。
        """
        started_at = time.time()
        presentation = Presentation(self.file_path)
        total_pages = len(presentation.slides)

        parsed_pages = self.parse_pptx_pages()
        return build_document_schema(
            file_path=self.file_path,
            document_type="pptx",
            total_pages=total_pages,
            parsed_pages=parsed_pages,
            started_at=started_at,
        )

    def parse(self) -> dict[str, Any]:
        """
        统一解析入口：根据文件扩展名自动选择解析方式。
        支持: .pdf / .pptx
        """
        suffix = Path(self.file_path).suffix.lower()
        if suffix == ".pdf":
            return self.parse_pdf()
        if suffix == ".pptx":
            return self.parse_pptx()
        raise ValueError(f"暂不支持的文件类型: {suffix}，目前仅支持 .pdf 和 .pptx")

    @staticmethod
    def dump_json(data: dict[str, Any], output_path: str | None = None) -> str:
        """将结构化结果转为JSON字符串，并可选写入文件。"""
        result = json.dumps(data, ensure_ascii=False, indent=2)
        if output_path:
            with open(output_path, "w", encoding="utf-8") as file:
                file.write(result)
        return result

    def _extract_pdf_text(self, page: fitz.Page) -> str:
        blocks = page.get_text("blocks") or []
        ordered_lines: list[str] = []

        for block in sorted(blocks, key=lambda item: (round(item[1], 1), round(item[0], 1))):
            block_text = str(block[4] or "").strip()
            if not block_text:
                continue
            for raw_line in block_text.splitlines():
                normalized = self._normalize_line(raw_line)
                if normalized:
                    ordered_lines.append(normalized)

        if ordered_lines:
            return "\n".join(self._merge_lines(ordered_lines))

        fallback = page.get_text("text") or ""
        return fallback.strip()

    @staticmethod
    def _normalize_line(line: str) -> str:
        return re.sub(r"[ \t\u3000]+", " ", line.strip())

    @staticmethod
    def _merge_lines(lines: list[str]) -> list[str]:
        merged: list[str] = []
        for line in lines:
            if not merged:
                merged.append(line)
                continue

            previous = merged[-1]
            if DocumentParser._should_merge(previous, line):
                merged[-1] = f"{previous} {line}".strip()
                continue

            merged.append(line)
        return merged

    @staticmethod
    def _should_merge(previous: str, current: str) -> bool:
        bullet_prefixes = ("-", "*", "•", "1.", "2.", "3.", "4.", "5.")
        if previous.endswith((":", "：", "?", "？", "!", "！")):
            return False
        if current.startswith(bullet_prefixes):
            return False
        if len(previous) <= 10:
            return False
        return not re.match(r"^[一二三四五六七八九十]+[、.]", current)

    @staticmethod
    def _is_meaningful_text(text: str) -> bool:
        compact = re.sub(r"\s+", "", text or "")
        return len(compact) >= 2 and bool(re.search(r"[A-Za-z0-9\u4e00-\u9fff]", compact))

    def _clean_text(self, text: str) -> str:
        """保留段落边界，减少解析后结构信息丢失。"""
        normalized = text.replace("\r\n", "\n").replace("\r", "\n")
        cleaned_lines: list[str] = []
        blank_pending = False

        for raw_line in normalized.split("\n"):
            line = self._normalize_line(raw_line)
            if not line:
                if cleaned_lines:
                    blank_pending = True
                continue

            if blank_pending and cleaned_lines:
                cleaned_lines.append("")
            cleaned_lines.append(line)
            blank_pending = False

        return "\n".join(cleaned_lines).strip()


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="测试 DocumentParser")
    parser.add_argument("file", help="要解析的 PDF 文件路径（示例：test.pdf）")
    parser.add_argument("--out", help="可选：输出JSON文件路径", default=None)
    args = parser.parse_args()

    dp = DocumentParser(args.file)
    result = dp.parse()
    if args.out:
        dp.dump_json(result, args.out)
        print(f"解析完成，结果已保存: {args.out}")
    else:
        print(dp.dump_json(result))
