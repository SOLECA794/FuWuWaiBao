import json
from pathlib import Path

import fitz
from pptx import Presentation

try:
    from .parser import DocumentParser
    from .generator import LessonGenerator, GenerationConfig
    from .qa import QAResponder, QAConfig
    from .reconstructor import LessonReconstructor
except ImportError:
    from parser import DocumentParser
    from generator import LessonGenerator, GenerationConfig
    from qa import QAResponder, QAConfig
    from reconstructor import LessonReconstructor


def _ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def _create_test_pdf(path: Path) -> None:
    doc = fitz.open()
    page1 = doc.new_page()
    page1.insert_text((72, 72), "Linear regression models the relation between variables")
    page2 = doc.new_page()
    page2.insert_text((72, 72), "Loss functions measure prediction error")
    doc.save(str(path))
    doc.close()


def _create_test_pptx(path: Path) -> None:
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "机器学习基础"
    slide1.placeholders[1].text = "监督学习与非监督学习"
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "模型评估"
    slide2.placeholders[1].text = "准确率、召回率、F1"
    prs.save(str(path))


def _save_json(path: Path, data: dict) -> None:
    with open(path, "w", encoding="utf-8") as file:
        json.dump(data, file, ensure_ascii=False, indent=2)


def _check(condition: bool, message: str, checks: list[dict]) -> None:
    checks.append({"check": message, "passed": bool(condition)})


def run_debug() -> dict:
    root = Path(__file__).resolve().parent
    out_dir = root / "debug_output"
    _ensure_dir(out_dir)

    checks: list[dict] = []

    pdf_path = out_dir / "sample.pdf"
    pptx_path = out_dir / "sample.pptx"
    _create_test_pdf(pdf_path)
    _create_test_pptx(pptx_path)

    # 1) 文档解析验证
    pdf_parsed = DocumentParser(str(pdf_path)).parse()
    pptx_parsed = DocumentParser(str(pptx_path)).parse()
    _save_json(out_dir / "parsed_pdf.json", pdf_parsed)
    _save_json(out_dir / "parsed_pptx.json", pptx_parsed)

    _check(pdf_parsed.get("doc_type") == "pdf", "PDF doc_type 正确", checks)
    _check(pptx_parsed.get("doc_type") == "pptx", "PPTX doc_type 正确", checks)
    _check(len(pdf_parsed.get("parsed_pages", [])) >= 1, "PDF 至少解析出 1 页", checks)
    _check(len(pptx_parsed.get("parsed_pages", [])) >= 1, "PPTX 至少解析出 1 页", checks)

    # 2) 讲稿与导图生成验证（llm）
    generator = LessonGenerator(GenerationConfig(mode="llm", model="qwen-plus", temperature=0.2))
    generated_pdf = generator.generate(pdf_parsed)
    _save_json(out_dir / "generated_pdf_mock.json", generated_pdf)

    reconstructed = LessonReconstructor().reconstruct(pdf_parsed)
    _save_json(out_dir / "reconstructed_pdf.json", reconstructed)

    lessons = generated_pdf.get("lessons", [])
    _check(len(lessons) >= 1, "生成结果包含至少 1 条 lesson", checks)
    if lessons:
        _check(bool(lessons[0].get("script")), "lesson 包含 script", checks)
        _check(bool(lessons[0].get("mindmap_markdown")), "lesson 包含 mindmap_markdown", checks)

    nodes = reconstructed.get("teaching_nodes", [])
    _check(len(nodes) >= 1, "重构结果包含至少 1 个 teaching_node", checks)
    if nodes:
        node_script = generator.generate_node_script(nodes[0], pdf_parsed.get("doc_name"))
        _save_json(out_dir / "generated_node.json", node_script)
        _check(bool(node_script.get("script")), "节点讲稿包含 script", checks)
        _check(bool(node_script.get("interactive_questions")), "节点讲稿包含互动问题", checks)

    # 3) 问答溯源 + 重讲验证（llm）
    responder = QAResponder(parsed_document=pdf_parsed, config=QAConfig(mode="llm", model="qwen-plus", temperature=0.2))
    qa_normal = responder.answer(question="这页主要在讲什么", current_page=1)
    qa_reteach = responder.answer(question="我听不懂，换个例子", current_page=1)
    _save_json(out_dir / "qa_normal.json", qa_normal)
    _save_json(out_dir / "qa_reteach.json", qa_reteach)

    _check(qa_normal.get("source_page") == 1, "普通问答返回正确 source_page", checks)
    _check(qa_normal.get("intent", {}).get("need_reteach") is False, "普通问答不触发重讲", checks)
    _check(qa_reteach.get("intent", {}).get("need_reteach") is True, "听不懂问题触发重讲", checks)
    _check(bool(qa_reteach.get("answer")), "重讲回答非空", checks)

    passed = all(item["passed"] for item in checks)
    summary = {
        "passed": passed,
        "total_checks": len(checks),
        "passed_checks": sum(1 for item in checks if item["passed"]),
        "checks": checks,
        "artifacts_dir": str(out_dir),
    }
    _save_json(out_dir / "summary.json", summary)
    return summary


if __name__ == "__main__":
    result = run_debug()
    print(json.dumps(result, ensure_ascii=False, indent=2))
